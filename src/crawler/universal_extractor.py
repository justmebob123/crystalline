#!/usr/bin/env python3
"""
Universal Text Extractor for Crawler
Supports: DOCX, XLSX, PPTX, ODT, ODS, ODP, EPUB, MD, YAML, TOML, SQL, and more
"""

import sys
import os
import json
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

# Office formats
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# eBook formats
try:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    EPUB_AVAILABLE = True
except ImportError:
    EPUB_AVAILABLE = False

# Structured data
try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False


def extract_docx(filepath):
    """Extract text from DOCX file"""
    if not DOCX_AVAILABLE:
        return None
    
    try:
        doc = Document(filepath)
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        return '\n\n'.join(text_parts)
    except Exception as e:
        print(f"Error extracting DOCX: {e}", file=sys.stderr)
        return None


def extract_xlsx(filepath):
    """Extract text from XLSX file"""
    if not XLSX_AVAILABLE:
        return None
    
    try:
        wb = load_workbook(filepath, read_only=True, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    text_parts.append(row_text)
        
        wb.close()
        return '\n'.join(text_parts)
    except Exception as e:
        print(f"Error extracting XLSX: {e}", file=sys.stderr)
        return None


def extract_pptx(filepath):
    """Extract text from PPTX file"""
    if not PPTX_AVAILABLE:
        return None
    
    try:
        prs = Presentation(filepath)
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n=== Slide {i} ===\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
                
                # Extract table content
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            text_parts.append(row_text)
        
        return '\n\n'.join(text_parts)
    except Exception as e:
        print(f"Error extracting PPTX: {e}", file=sys.stderr)
        return None


def extract_odt(filepath):
    """Extract text from ODT file using odt2txt"""
    try:
        import subprocess
        result = subprocess.run(['odt2txt', filepath], 
                              capture_output=True, text=True, timeout=30)
        if result.returncode == 0:
            return result.stdout
        return None
    except Exception as e:
        print(f"Error extracting ODT: {e}", file=sys.stderr)
        return None


def extract_ods(filepath):
    """Extract text from ODS file"""
    try:
        # ODS is a zip file with XML content
        with zipfile.ZipFile(filepath, 'r') as zip_ref:
            content_xml = zip_ref.read('content.xml')
            root = ET.fromstring(content_xml)
            
            # Extract all text nodes
            text_parts = []
            for text_elem in root.iter():
                if text_elem.text and text_elem.text.strip():
                    text_parts.append(text_elem.text.strip())
            
            return '\n'.join(text_parts)
    except Exception as e:
        print(f"Error extracting ODS: {e}", file=sys.stderr)
        return None


def extract_odp(filepath):
    """Extract text from ODP file"""
    try:
        # ODP is a zip file with XML content
        with zipfile.ZipFile(filepath, 'r') as zip_ref:
            content_xml = zip_ref.read('content.xml')
            root = ET.fromstring(content_xml)
            
            # Extract all text nodes
            text_parts = []
            for text_elem in root.iter():
                if text_elem.text and text_elem.text.strip():
                    text_parts.append(text_elem.text.strip())
            
            return '\n'.join(text_parts)
    except Exception as e:
        print(f"Error extracting ODP: {e}", file=sys.stderr)
        return None


def extract_epub(filepath):
    """Extract text from EPUB file"""
    if not EPUB_AVAILABLE:
        return None
    
    try:
        book = epub.read_epub(filepath)
        text_parts = []
        
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text = soup.get_text(separator='\n', strip=True)
                if text:
                    text_parts.append(text)
        
        return '\n\n'.join(text_parts)
    except Exception as e:
        print(f"Error extracting EPUB: {e}", file=sys.stderr)
        return None


def extract_markdown(filepath):
    """Extract text from Markdown file (preserve structure)"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Preserve markdown structure but clean up
        lines = []
        for line in content.split('\n'):
            # Keep headers, lists, code blocks
            if line.strip():
                lines.append(line)
        
        return '\n'.join(lines)
    except Exception as e:
        print(f"Error extracting Markdown: {e}", file=sys.stderr)
        return None


def extract_yaml(filepath):
    """Extract text from YAML file"""
    if not YAML_AVAILABLE:
        return None
    
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            data = yaml.safe_load(f)
        
        # Convert YAML to readable text
        def yaml_to_text(obj, indent=0):
            text_parts = []
            prefix = '  ' * indent
            
            if isinstance(obj, dict):
                for key, value in obj.items():
                    text_parts.append(f"{prefix}{key}:")
                    text_parts.append(yaml_to_text(value, indent + 1))
            elif isinstance(obj, list):
                for item in obj:
                    text_parts.append(yaml_to_text(item, indent))
            else:
                text_parts.append(f"{prefix}{obj}")
            
            return '\n'.join(text_parts)
        
        return yaml_to_text(data)
    except Exception as e:
        print(f"Error extracting YAML: {e}", file=sys.stderr)
        return None


def extract_toml(filepath):
    """Extract text from TOML file"""
    try:
        # TOML is similar to INI format, just read as text
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"Error extracting TOML: {e}", file=sys.stderr)
        return None


def extract_sql(filepath):
    """Extract text from SQL file"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Clean up SQL comments
        lines = []
        for line in content.split('\n'):
            line = line.strip()
            if line and not line.startswith('--'):
                lines.append(line)
        
        return '\n'.join(lines)
    except Exception as e:
        print(f"Error extracting SQL: {e}", file=sys.stderr)
        return None


def extract_latex(filepath):
    """Extract text from LaTeX file"""
    try:
        import subprocess
        # Use pandoc to convert LaTeX to plain text
        result = subprocess.run(['pandoc', '-f', 'latex', '-t', 'plain', filepath],
                              capture_output=True, text=True, timeout=30)
        if result.returncode == 0:
            return result.stdout
        
        # Fallback: just read the file
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"Error extracting LaTeX: {e}", file=sys.stderr)
        return None


def extract_eml(filepath):
    """Extract text from EML (email) file"""
    try:
        import email
        from email import policy
        
        with open(filepath, 'rb') as f:
            msg = email.message_from_binary_file(f, policy=policy.default)
        
        text_parts = []
        
        # Extract headers
        text_parts.append(f"From: {msg.get('From', 'Unknown')}")
        text_parts.append(f"To: {msg.get('To', 'Unknown')}")
        text_parts.append(f"Subject: {msg.get('Subject', 'No Subject')}")
        text_parts.append(f"Date: {msg.get('Date', 'Unknown')}")
        text_parts.append("\n")
        
        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    text_parts.append(part.get_content())
        else:
            text_parts.append(msg.get_content())
        
        return '\n'.join(text_parts)
    except Exception as e:
        print(f"Error extracting EML: {e}", file=sys.stderr)
        return None


def detect_and_extract(filepath):
    """Detect file type and extract text"""
    ext = Path(filepath).suffix.lower()
    
    extractors = {
        '.docx': extract_docx,
        '.xlsx': extract_xlsx,
        '.pptx': extract_pptx,
        '.odt': extract_odt,
        '.ods': extract_ods,
        '.odp': extract_odp,
        '.epub': extract_epub,
        '.md': extract_markdown,
        '.markdown': extract_markdown,
        '.yaml': extract_yaml,
        '.yml': extract_yaml,
        '.toml': extract_toml,
        '.sql': extract_sql,
        '.tex': extract_latex,
        '.eml': extract_eml,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        return extractor(filepath)
    
    return None


def main():
    if len(sys.argv) != 2:
        print("Usage: universal_extractor.py <filepath>", file=sys.stderr)
        sys.exit(1)
    
    filepath = sys.argv[1]
    
    if not os.path.exists(filepath):
        print(f"File not found: {filepath}", file=sys.stderr)
        sys.exit(1)
    
    text = detect_and_extract(filepath)
    
    if text:
        print(text)
        sys.exit(0)
    else:
        print(f"Could not extract text from: {filepath}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()